import os
import re
import tempfile
import fitz  # PyMuPDF
from pptx import Presentation
from fastapi import FastAPI, Request, BackgroundTasks
from google import genai
from supabase import create_client, Client

app = FastAPI()

# Environment variables
SUPABASE_URL = os.environ.get("SUPABASE_URL", "")
SUPABASE_SERVICE_KEY = os.environ.get("SUPABASE_SERVICE_KEY", "")
GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY", "")

if SUPABASE_URL and SUPABASE_SERVICE_KEY:
    supabase: Client = create_client(SUPABASE_URL, SUPABASE_SERVICE_KEY)
if GEMINI_API_KEY:
    gemini_client = genai.Client(api_key=GEMINI_API_KEY)

def extract_text_from_pdf(filepath):
    doc = fitz.open(filepath)
    text = ""
    for page in doc:
        text += page.get_text() + "\n"
    return text

def extract_text_from_pptx(filepath):
    prs = Presentation(filepath)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def chunk_text(text, chunk_size=800, overlap=150):
    words = text.split()
    chunks = []
    i = 0
    while i < len(words):
        chunk = " ".join(words[i:i + chunk_size])
        chunks.append(chunk)
        i += chunk_size - overlap
    return chunks

def process_file_background(file_path: str):
    print(f"Processing {file_path}...")
    
    # 1. Download file from Supabase
    bucket = "course_files"
    res = supabase.storage.from_(bucket).download(file_path)
    
    # Create temp file
    ext = os.path.splitext(file_path)[1].lower()
    with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as temp_file:
        temp_file.write(res)
        temp_file_path = temp_file.name

    try:
        # 2. Extract text
        text = ""
        if ext == ".pdf":
            text = extract_text_from_pdf(temp_file_path)
        elif ext == ".pptx":
            text = extract_text_from_pptx(temp_file_path)
        else:
            print(f"Unsupported file type: {ext}")
            return
            
        if not text.strip():
            print("No text found.")
            return

        # 3. Chunk text
        chunks = chunk_text(text)
        print(f"Split into {len(chunks)} chunks.")

        # Parse course_id from filepath (course_{course_id}_mod_{mod_id}_{filename})
        match = re.search(r'course_(\d+)_', file_path)
        course_id = match.group(1) if match else "unknown"

        # 4. Generate embeddings and insert
        for i, chunk in enumerate(chunks):
            try:
                embedding_res = gemini_client.models.embed_content(
                    model="text-embedding-004",
                    contents=chunk
                )
                embedding = embedding_res.embeddings[0].values
                
                supabase.table('course_embeddings').insert({
                    "course_id": course_id,
                    "file_path": file_path,
                    "chunk_index": i,
                    "chunk_text": chunk,
                    "embedding": embedding
                }).execute()
            except Exception as e:
                print(f"Failed to embed chunk {i}: {e}")

        print(f"Successfully indexed {file_path}")

    except Exception as e:
        print(f"Error processing {file_path}: {e}")
    finally:
        os.remove(temp_file_path)

@app.post("/webhook/supabase")
async def supabase_webhook(request: Request, background_tasks: BackgroundTasks):
    """
    Receives Database Webhook from Supabase when a new file is added to storage.objects.
    """
    try:
        payload = await request.json()
    except:
        return {"status": "error", "message": "Invalid JSON"}
    
    # Check if this is an INSERT into the storage.objects table
    if payload.get("type") == "INSERT" and payload.get("table") == "objects":
        record = payload.get("record", {})
        bucket_id = record.get("bucket_id")
        file_path = record.get("name")
        
        if bucket_id == "course_files" and file_path:
            # Run the heavy processing in the background so the webhook responds quickly to Supabase
            background_tasks.add_task(process_file_background, file_path)
            return {"status": "processing_started", "file": file_path}
            
    return {"status": "ignored"}

@app.get("/")
def health_check():
    return {"status": "NotMoodle RAG Backend is running!"}
